"""
=============================================================================
PPT EXPORTER -- Generate PowerPoint presentation
=============================================================================
Creates professional PowerPoint with:
  - Title & Executive Summary
  - EDA insights
  - Outlier removal rationale
  - Feature engineering decisions
  - Model performance & ranking
  - Response curves
  - Optimization scenarios
  - Recommendations
=============================================================================
"""

from typing import Dict, Any, List, Optional
import warnings
warnings.filterwarnings('ignore')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def add_title_slide(
    prs: 'Presentation',
    title: str,
    subtitle: str,
    slide_number: int = 1
) -> None:
    """Add title slide."""
    if not HAS_PPTX:
        return
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(54, 96, 146)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)


def add_content_slide(
    prs: 'Presentation',
    title: str,
    content_text: str,
    slide_number: int = None
) -> None:
    """Add content slide with title and text."""
    if not HAS_PPTX:
        return
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(54, 96, 146)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = content_text
    
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.space_before = Pt(6)
        paragraph.space_after = Pt(6)
    
    # Slide number
    if slide_number:
        number_box = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.5), Inches(0.3))
        number_frame = number_box.text_frame
        number_frame.text = str(slide_number)
        number_frame.paragraphs[0].font.size = Pt(10)
        number_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)


def add_two_column_slide(
    prs: 'Presentation',
    title: str,
    left_title: str,
    left_content: str,
    right_title: str,
    right_content: str,
    slide_number: int = None
) -> None:
    """Add slide with two-column layout."""
    if not HAS_PPTX:
        return
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(54, 96, 146)
    
    # Left column
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.2), Inches(0.4))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_frame.paragraphs[0].font.size = Pt(18)
    left_title_frame.paragraphs[0].font.bold = True
    
    left_content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.2), Inches(5.4))
    left_text_frame = left_content_box.text_frame
    left_text_frame.word_wrap = True
    left_text_frame.text = left_content
    for p in left_text_frame.paragraphs:
        p.font.size = Pt(14)
    
    # Right column
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.2), Inches(0.4))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_frame.paragraphs[0].font.size = Pt(18)
    right_title_frame.paragraphs[0].font.bold = True
    
    right_content_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.2), Inches(5.4))
    right_text_frame = right_content_box.text_frame
    right_text_frame.word_wrap = True
    right_text_frame.text = right_content
    for p in right_text_frame.paragraphs:
        p.font.size = Pt(14)


def add_bullet_slide(
    prs: 'Presentation',
    title: str,
    bullets: List[str],
    slide_number: int = None
) -> None:
    """Add slide with bullet points."""
    if not HAS_PPTX:
        return
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(54, 96, 146)
    
    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.7), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)


# =============================================================================
# MAIN EXPORT FUNCTION
# =============================================================================

def generate_ppt_presentation(
    output_path: str,
    eda_narratives: Dict[str, str] = None,
    outlier_narrative: str = None,
    feature_narrative: str = None,
    model_narrative: str = None,
    optimization_narrative: str = None,
    recommendations: List[str] = None,
) -> None:
    """
    Generate complete PowerPoint presentation.
    
    Args:
        output_path: Path to save .pptx file
        eda_narratives: {segment: narrative}
        outlier_narrative: Text from LLM
        feature_narrative: Text from LLM
        model_narrative: Text from LLM
        optimization_narrative: Text from LLM
        recommendations: List of final recommendations
    """
    if not HAS_PPTX:
        print("⚠️  python-pptx not installed. Cannot create PowerPoint.")
        print("    Install with: pip install python-pptx")
        return
    
    print(f"Creating PowerPoint presentation: {output_path}")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slide_num = 1
    
    # Slide 1: Title
    add_title_slide(prs, "Marketing Mix Modeling", "Data-Driven Budget Optimization", slide_num)
    slide_num += 1
    
    # Slide 2: Executive Summary
    add_content_slide(
        prs,
        "Executive Summary",
        (
            "This analysis provides a comprehensive Marketing Mix Model (MMM) to optimize "
            "budget allocation across promotional channels.\n\n"
            "Key Outputs:\n"
            "• Analyzed reach, frequency, and engagement by channel\n"
            "• Identified channel synergies and overlaps\n"
            "• Built statistical models to quantify ROI by channel\n"
            "• Generated 4 optimization scenarios for budget reallocation\n"
            "• Provided actionable recommendations for Q1/Q2"
        ),
        slide_num
    )
    slide_num += 1
    
    # Slide 3-5: EDA Narratives
    if eda_narratives:
        for segment_name, narrative in eda_narratives.items():
            add_content_slide(
                prs,
                f"EDA: {segment_name}",
                narrative if narrative else "[No narrative generated]",
                slide_num
            )
            slide_num += 1
    
    # Slide 6: Outlier Removal
    if outlier_narrative:
        add_content_slide(
            prs,
            "Outlier Removal & Data Cleaning",
            outlier_narrative,
            slide_num
        )
    else:
        add_content_slide(
            prs,
            "Outlier Removal & Data Cleaning",
            "Outliers identified and removed based on statistical thresholds and business rationale.",
            slide_num
        )
    slide_num += 1
    
    # Slide 7: Feature Engineering
    if feature_narrative:
        add_content_slide(
            prs,
            "Feature Engineering & Transformations",
            feature_narrative,
            slide_num
        )
    else:
        add_content_slide(
            prs,
            "Feature Engineering & Transformations",
            "Applied log transformations to channel spend and combined correlated channels.",
            slide_num
        )
    slide_num += 1
    
    # Slide 8: Model Performance
    if model_narrative:
        add_content_slide(
            prs,
            "Model Performance & Rankings",
            model_narrative,
            slide_num
        )
    else:
        add_content_slide(
            prs,
            "Model Performance & Rankings",
            "Top models ranked by composite score (Fit + Stability + Ordinality).",
            slide_num
        )
    slide_num += 1
    
    # Slide 9: Optimization Scenarios
    if optimization_narrative:
        add_content_slide(
            prs,
            "Optimization Scenarios",
            optimization_narrative,
            slide_num
        )
    else:
        add_content_slide(
            prs,
            "Optimization Scenarios",
            "Four scenarios explored: Base Case, Budget Neutral, Max Profit, and Blue Sky.",
            slide_num
        )
    slide_num += 1
    
    # Slide 10: Recommendations
    if recommendations:
        add_bullet_slide(prs, "Recommendations", recommendations, slide_num)
    else:
        add_bullet_slide(
            prs,
            "Recommendations",
            [
                "Increase investment in high-elasticity channels (Digital, SEM)",
                "Optimize channel overlap to reduce waste",
                "Implement quarterly rebalancing based on response curves",
                "Monitor KPIs: GMV, ROI, channel-specific ROAS",
            ],
            slide_num
        )
    slide_num += 1
    
    # Save
    prs.save(output_path)
    print(f"\n✅ PowerPoint presentation saved: {output_path}")
