"""
=============================================================================
POWERPOINT EXPORTER (THIN WRAPPER)
=============================================================================
Thin wrapper that creates PowerPoint presentations from pipeline results.

NOTE: Can integrate deeper with specific slide layouts later.
Currently focuses on summary + key metrics.
"""

from typing import Dict, Any, List, Optional
import warnings
warnings.filterwarnings('ignore')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# =============================================================================
# UTILITIES
# =============================================================================

def add_title_slide(prs, title: str, subtitle: str = ""):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(54, 96, 146)
    
    # Subtitle
    if subtitle:
        left = Inches(0.5)
        top = Inches(3.8)
        width = Inches(9)
        height = Inches(1)
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)


def add_content_slide(prs, title: str, content_dict: Dict[str, Any]):
    """Add content slide with title and key metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(54, 96, 146)
    
    # Content
    left = Inches(0.75)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(4.5)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for key, value in content_dict.items():
        p = content_frame.add_paragraph()
        p.text = f"{key}: {value}"
        p.font.size = Pt(16)
        p.level = 0
        p.space_after = Pt(8)


def add_metrics_slide(prs, title: str, metrics: Dict[str, float]):
    """Add slide with formatted metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(54, 96, 146)
    
    # Metrics in 2 columns
    left = Inches(0.75)
    top = Inches(1.7)
    width = Inches(4)
    height = Inches(3.5)
    
    col1_box = slide.shapes.add_textbox(left, top, width, height)
    col1_frame = col1_box.text_frame
    col1_frame.word_wrap = True
    
    col2_box = slide.shapes.add_textbox(left + Inches(4.5), top, width, height)
    col2_frame = col2_box.text_frame
    col2_frame.word_wrap = True
    
    items = list(metrics.items())
    mid = (len(items) + 1) // 2
    
    for key, value in items[:mid]:
        p = col1_frame.add_paragraph()
        
        # Format value
        if isinstance(value, float):
            if abs(value) < 1:
                p.text = f"{key}: {value:.4f}"
            else:
                p.text = f"{key}: {value:.2f}"
        else:
            p.text = f"{key}: {value}"
        
        p.font.size = Pt(14)
        p.space_after = Pt(10)
    
    for key, value in items[mid:]:
        p = col2_frame.add_paragraph()
        
        if isinstance(value, float):
            if abs(value) < 1:
                p.text = f"{key}: {value:.4f}"
            else:
                p.text = f"{key}: {value:.2f}"
        else:
            p.text = f"{key}: {value}"
        
        p.font.size = Pt(14)
        p.space_after = Pt(10)


# =============================================================================
# EXPORT TO POWERPOINT
# =============================================================================

def export_to_powerpoint(
    pipeline_results: Dict[str, Any],
    output_path: str
) -> Dict[str, Any]:
    """
    Export pipeline results to PowerPoint presentation.
    
    Args:
        pipeline_results: Dict from orchestrator with all results
        output_path: Path to write .pptx file
        
    Returns:
        {'success': bool, 'path': str, 'message': str}
    """
    if not HAS_PPTX:
        return {
            'success': False,
            'message': 'python-pptx not installed. Install: pip install python-pptx'
        }
    
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # =====================================================================
        # SLIDE 1: TITLE
        # =====================================================================
        add_title_slide(
            prs,
            "Marketing Mix Analysis",
            "Optimization Results & Recommendations"
        )
        
        # =====================================================================
        # SLIDE 2: SUMMARY
        # =====================================================================
        if 'metadata' in pipeline_results:
            meta = pipeline_results['metadata']
            add_content_slide(prs, "Pipeline Summary", {
                'Mode': meta.get('mode', 'N/A'),
                'Data Rows': meta.get('data_rows', 'N/A'),
                'Channels': meta.get('channels', 'N/A'),
            })
        
        # =====================================================================
        # SLIDE 3: EDA RESULTS
        # =====================================================================
        if 'eda_metrics' in pipeline_results:
            eda = pipeline_results['eda_metrics']
            
            if 'eda_results' in eda and 'National' in eda['eda_results']:
                nat = eda['eda_results']['National']
                
                if 'rfe' in nat:
                    rfe = nat['rfe']
                    add_metrics_slide(prs, "EDA: Reach & Frequency", {
                        'Reach': rfe.get('reach', 0),
                        'Frequency': rfe.get('frequency', 0),
                        'Engagement': rfe.get('engagement', 0),
                    })
        
        # =====================================================================
        # SLIDE 4: BEST MODEL
        # =====================================================================
        if 'modeling' in pipeline_results:
            mod = pipeline_results['modeling']
            
            if 'best_model' in mod and mod['best_model']:
                best = mod['best_model']
                
                # Handle both dict and object
                if hasattr(best, '__dict__'):
                    best_dict = best.__dict__
                else:
                    best_dict = best
                
                add_metrics_slide(prs, "Best Performing Model", {
                    'Type': best_dict.get('model_type', 'N/A'),
                    'R² Score': best_dict.get('r2', 0),
                    'RMSE': best_dict.get('rmse', 0),
                    'Features': best_dict.get('feature_count', 0),
                    'Overall Score': best_dict.get('overall_score', 0),
                })
        
        # =====================================================================
        # SLIDE 5: OPTIMIZATION RECOMMENDATIONS
        # =====================================================================
        if 'optimization' in pipeline_results:
            opt = pipeline_results['optimization']
            
            if 'scenarios' in opt and 'max_profit' in opt['scenarios']:
                max_profit = opt['scenarios']['max_profit']
                
                if max_profit and 'error' not in max_profit:
                    add_metrics_slide(prs, "Max Profit Scenario", {
                        'Total Spend': f"${max_profit.get('total_spend', 0):.0f}",
                        'Predicted Sales': f"${max_profit.get('predicted_sales', 0):.0f}",
                        'Profit': f"${max_profit.get('profit', 0):.0f}",
                        'ROI': f"{max_profit.get('roi', 0):.2%}",
                    })
        
        # =====================================================================
        # SLIDE 6: KEY INSIGHTS
        # =====================================================================
        insights = []
        
        if 'eda_metrics' in pipeline_results and pipeline_results['eda_metrics'].get('is_valid'):
            insights.append("✓ EDA: Data quality validated")
        
        if 'modeling' in pipeline_results and pipeline_results['modeling'].get('is_valid'):
            mod = pipeline_results['modeling']
            if mod['best_model']:
                best = mod['best_model']
                r2 = best.get('r2') if isinstance(best, dict) else best.r2
                insights.append(f"✓ Best model fit: R² = {r2:.3f}")
        
        if 'optimization' in pipeline_results and pipeline_results['optimization'].get('is_valid'):
            insights.append("✓ Optimization complete: 4 scenarios analyzed")
        
        if insights:
            add_content_slide(prs, "Key Insights", {
                f"Insight {i+1}": insight for i, insight in enumerate(insights)
            })
        
        # Save
        prs.save(output_path)
        
        return {
            'success': True,
            'path': output_path,
            'message': f'PowerPoint exported: {output_path}',
            'slides': len(prs.slides),
        }
    
    except Exception as e:
        return {
            'success': False,
            'path': output_path,
            'message': f'Export failed: {str(e)}',
        }
